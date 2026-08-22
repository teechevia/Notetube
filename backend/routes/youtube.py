from fastapi import APIRouter
from models.request_models import VideoRequest, TranscriptRequest
from services.transcript import process_video
from services.generators import generate_flashcards, generate_quiz, generate_ppt, generate_podcast

router = APIRouter()

@router.get("/")
def home():
    return {"message": "LectureLens AI backend running"}
 
@router.post("/video")
def receive_video(video: VideoRequest):
    return process_video(video.url)

@router.post("/flashcards")
def get_flashcards(req: TranscriptRequest):
    return {"flashcards": generate_flashcards(req.transcript)}

@router.post("/quiz")
def get_quiz(req: TranscriptRequest):
    return {"quiz": generate_quiz(req.transcript)}

@router.post("/ppt")
def get_ppt(req: TranscriptRequest):
    return {"ppt": generate_ppt(req.transcript)}

@router.post("/podcast")
def get_podcast(req: TranscriptRequest):
    return {"podcast": generate_podcast(req.transcript)}

from fastapi.responses import FileResponse
from models.request_models import TextRequest
try:
    from pptx import Presentation
except ImportError:
    pass
try:
    from gtts import gTTS
except ImportError:
    pass

@router.post("/ppt/download")
def download_ppt(req: TextRequest):
    prs = Presentation()
    slides_text = req.text.split("## Slide:")
    for slide_text in slides_text:
        if not slide_text.strip(): continue
        lines = slide_text.strip().split('\n')
        title = lines[0].strip()
        bullets = [l.replace('- ', '').replace('* ', '').strip() for l in lines[1:] if l.strip()]
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        tf = slide.placeholders[1].text_frame
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
            
    filename = "presentation.pptx"
    prs.save(filename)
    return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename="presentation.pptx")

@router.post("/podcast/download")
def download_podcast(req: TextRequest):
    tts = gTTS(req.text)
    filename = "podcast.mp3"
    tts.save(filename)
    return FileResponse(filename, media_type="audio/mpeg", filename="podcast.mp3")
