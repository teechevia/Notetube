
from langchain_community.document_loaders import WebBaseLoader, PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from youtube_transcript_api import YouTubeTranscriptApi
import re
import os

text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

def extract_video_id(url):
    match = re.search(r"(?:v=|\/)([0-9A-Za-z_-]{11}).*", url)
    return match.group(1) if match else None

def load_youtube(url):
    try:
        video_id = extract_video_id(url)
        if not video_id:
            raise ValueError("Invalid YouTube URL")

        api = YouTubeTranscriptApi()
        tx_list = api.list(video_id)
        try:
            tx = tx_list.find_transcript(['en'])
        except:
            try:
                tx = tx_list.find_generated_transcript(['en'])
            except:
                # Fallback to the first available transcript if English is not available
                tx = list(tx_list)[0]
        transcript = tx.fetch()

        full_text = []
        for t in transcript:
            # Format: [MM:SS] Text
            minutes = int(t.start // 60)
            seconds = int(t.start % 60)
            time_str = f"[{minutes:02d}:{seconds:02d}]"
            full_text.append(f"{time_str} {t.text}")

        final_text = "\n".join(full_text)

        doc = Document(page_content=final_text, metadata={"source": url, "type": "youtube"})
        return text_splitter.split_documents([doc])
    except Exception as e:
        print(f"Error loading YouTube: {e}")
        return []

def load_url(url):
    try:
        loader = WebBaseLoader(url)
        docs = loader.load()
        for doc in docs:
            doc.metadata["type"] = "url"
        return text_splitter.split_documents(docs)
    except Exception as e:
        print(f"Error loading URL: {e}")
        return []

def load_pdf(file_path):
    try:
        loader = PyPDFLoader(file_path)
        docs = loader.load()
        for doc in docs:
            doc.metadata["type"] = "pdf"
        return text_splitter.split_documents(docs)
    except Exception as e:
        print(f"Error loading PDF: {e}")
        return []

def load_txt(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        doc = Document(page_content=text, metadata={"source": os.path.basename(file_path), "type": "txt"})
        return text_splitter.split_documents([doc])
    except Exception as e:
        print(f"Error loading TXT: {e}")
        return []

def load_docx(file_path):
    try:
        import docx
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        text = "\n".join(full_text)
        d = Document(page_content=text, metadata={"source": os.path.basename(file_path), "type": "docx"})
        return text_splitter.split_documents([d])
    except Exception as e:
        print(f"Error loading DOCX: {e}")
        return []

def load_pptx(file_path):
    try:
        import pptx
        prs = pptx.Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        text = "\n".join(full_text)
        d = Document(page_content=text, metadata={"source": os.path.basename(file_path), "type": "pptx"})
        return text_splitter.split_documents([d])
    except Exception as e:
        print(f"Error loading PPTX: {e}")
        return []

def load_audio(file_path):
    # This requires Gemini API for transcription.
    # We will upload the audio file to Gemini and ask for a complete transcript.
    try:
        import google.generativeai as genai
        import os
        from dotenv import load_dotenv

        load_dotenv()
        api_key = os.getenv("GEMINI_API_KEY")
        if not api_key:
            raise ValueError("GEMINI_API_KEY missing")
        genai.configure(api_key=api_key)

        audio_file = genai.upload_file(path=file_path)
        model = genai.GenerativeModel("gemini-1.5-flash")

        prompt = "Provide a highly accurate, word-for-word transcript of this audio file. Do not add any extra commentary, just the transcript."
        response = model.generate_content([prompt, audio_file])

        text = response.text

        # Clean up
        genai.delete_file(audio_file.name)

        d = Document(page_content=text, metadata={"source": os.path.basename(file_path), "type": "audio"})
        return text_splitter.split_documents([d])
    except Exception as e:
        print(f"Error loading audio: {e}")
        return []

